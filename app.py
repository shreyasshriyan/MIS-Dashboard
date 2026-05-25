#!/usr/bin/env python3
"""
Logistics Dashboard — Flask app with HTTP Basic Auth.
Designed for Render.com deployment.
Set RENDER_USERNAME and RENDER_PASSWORD env vars for auth.
"""

import base64
import io
import os
import re
import zipfile
from datetime import datetime, date
from functools import wraps
from pathlib import Path

from flask import Flask, jsonify, request, send_file, render_template_string
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from docx import Document
from docx.shared import Inches as DocInches, Pt as DocPt, RGBColor as DocRGB
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from fpdf import FPDF

HERE = Path(__file__).parent

app = Flask(__name__, static_folder=str(HERE), static_url_path='/static')

# ── Auth config ──────────────────────────────────────────────────────
AUTH_USER = os.environ.get('RENDER_USERNAME', 'admin')
AUTH_PASS = os.environ.get('RENDER_PASSWORD', 'changeme')

def require_auth(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        auth = request.authorization
        if not auth or auth.username != AUTH_USER or auth.password != AUTH_PASS:
            return (
                'Unauthorized',
                401,
                {'WWW-Authenticate': 'Basic realm="Logistics Dashboard"'},
            )
        return f(*args, **kwargs)
    return decorated

# ── Colours ──────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1E, 0x3A, 0x5F)
BLUE   = RGBColor(0x25, 0x63, 0xEB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x0F, 0x17, 0x2A)
GRAY   = RGBColor(0x64, 0x74, 0x8B)
LGRAY  = RGBColor(0x94, 0xA3, 0xB8)

def fmt_num(n):
    return f'{int(n):,}' if n == int(n) else f'{n:,.1f}'

def fmt_money(n):
    n = float(n or 0)
    if abs(n) >= 10_000_000:
        return f'INR {n/10_000_000:.1f} Cr'
    if abs(n) >= 100_000:
        return f'INR {n/100_000:.1f} L'
    return f'INR {int(n):,}'

def fmt_date(d):
    if isinstance(d, str) and d:
        try:
            return datetime.strptime(d[:10], '%Y-%m-%d').strftime('%d %b %Y')
        except ValueError:
            return d
    return '-'

def safe_text(t):
    t = str(t)
    reps = {'\u2014': '--', '\u2013': '-', '\u2018': "'", '\u2019': "'",
            '\u201c': '"', '\u201d': '"', '\u2026': '...', '\u2022': '-',
            '\u20b9': 'Rs.', '\u00a0': ' '}
    for old, new in reps.items():
        t = t.replace(old, new)
    return t.encode('latin-1', errors='replace').decode('latin-1')

# ── Routes ───────────────────────────────────────────────────────────

@app.route('/')
@require_auth
def index():
    return send_file(str(HERE / 'index.html'))

@app.route('/api/generate', methods=['POST', 'OPTIONS'])
def api_generate():
    files = request.files.getlist('files')
    fmt = request.form.get('format', 'ppt')
    if not files:
        return jsonify({'error': 'No CSV files uploaded'}), 400

    csv_texts = {}
    for f in files:
        csv_texts[f.filename] = f.read().decode('utf-8', errors='replace')

    datasets = []
    for fname, text in csv_texts.items():
        records = parse_csv(text)
        if records:
            client = ''
            for k in records[0]:
                kl = k.lower().strip()
                if kl in ('client', 'customer code', 'customer', 'network'):
                    client = str(records[0].get(k, '')).replace('`', '').strip().title()
                    break
            if not client:
                client = fname.replace('.csv', '').replace('_', ' ').replace('-', ' ').title()
            datasets.append({'records': records, 'client': client})

    if not datasets:
        return jsonify({'error': 'No valid records'}), 400

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zf:
        for ds in datasets:
            report = build_report(ds['records'])
            report['client'] = ds['client']
            fn = safe_filename(ds['client'])
            if fmt == 'ppt':
                zf.writestr(f'{fn}_Logistics_Report.pptx', generate_ppt(report))
            elif fmt == 'word':
                zf.writestr(f'{fn}_Logistics_Report.docx', generate_docx(report))
            else:
                zf.writestr(f'{fn}_Logistics_Report.pdf', generate_pdf(report))

        if len(datasets) >= 2:
            all_recs = [r for ds in datasets for r in ds['records']]
            merged = build_report(all_recs)
            merged['client'] = 'All Customers (Consolidated)'
            if fmt == 'ppt':
                zf.writestr('All_Customers_Consolidated_Logistics_Report.pptx', generate_ppt(merged))
            elif fmt == 'word':
                zf.writestr('All_Customers_Consolidated_Logistics_Report.docx', generate_docx(merged))
            else:
                zf.writestr('All_Customers_Consolidated_Logistics_Report.pdf', generate_pdf(merged))

    buf.seek(0)
    return send_file(buf, mimetype='application/zip', as_attachment=True,
                     download_name=f'Logistics_Reports_{date.today().isoformat()}.zip')

# ── CSV Parsing ──────────────────────────────────────────────────────

def parse_csv(text):
    text = text.lstrip('\ufeff')
    lines = text.split('\n')
    if not lines:
        return []
    header = [c.strip().strip('"').strip("'") for c in lines[0].split(',')]
    records = []
    for line in lines[1:]:
        line = line.strip()
        if not line:
            continue
        vals = []
        current = ''
        quoted = False
        for ch in line:
            if quoted:
                if ch == '"':
                    quoted = False
                else:
                    current += ch
            elif ch == '"':
                quoted = True
            elif ch == ',':
                vals.append(current.strip())
                current = ''
            else:
                current += ch
        vals.append(current.strip())
        while len(vals) < len(header):
            vals.append('')
        rec = dict(zip(header, vals[:len(header)]))
        records.append(rec)
    return records

def parse_num(v):
    try:
        return float(str(v or '').replace(',', '').replace('`', ''))
    except ValueError:
        return 0.0

def parse_dt(v):
    s = str(v or '').strip()
    m = re.match(r'(\d{4})-(\d{2})-(\d{2})', s)
    if m:
        return datetime(int(m[1]), int(m[2]), int(m[3]))
    m = re.match(r'(\d{2})-(\d{2})-(\d{4})', s)
    if m:
        return datetime(int(m[3]), int(m[2]), int(m[1]))
    return None

def guess_state(city):
    c = (city or '').lower().strip()
    if not c:
        return 'Unknown'
    m = {
        'bengaluru':'Karnataka','bangalore':'Karnataka','mysore':'Karnataka','mangalore':'Karnataka',
        'hubli':'Karnataka','bellary':'Karnataka','bijapur':'Karnataka','bidar':'Karnataka',
        'shimoga':'Karnataka','belgaum':'Karnataka','hosur':'Karnataka','channagiri':'Karnataka',
        'arsikere':'Karnataka','hassan':'Karnataka','kolar':'Karnataka',
        'mumbai':'Maharashtra','pune':'Maharashtra','nagpur':'Maharashtra','thane':'Maharashtra',
        'navi mumbai':'Maharashtra','greater thane':'Maharashtra','nashik':'Maharashtra',
        'aurangabad':'Maharashtra','kolhapur':'Maharashtra','solapur':'Maharashtra',
        'sangli':'Maharashtra','jalgaon':'Maharashtra','ahmednagar':'Maharashtra',
        'akola':'Maharashtra','chandrapur':'Maharashtra','gadchiroli':'Maharashtra',
        'osmanabad':'Maharashtra','ratnagiri':'Maharashtra','satara':'Maharashtra',
        'palghar':'Maharashtra','mangaon':'Maharashtra','nagothane':'Maharashtra',
        'khed':'Maharashtra','phali':'Maharashtra','phaltan':'Maharashtra',
        'shegaon':'Maharashtra','pusad':'Maharashtra','pen':'Maharashtra','alibag':'Maharashtra',
        'delhi':'Delhi','new delhi':'Delhi','delhi ncr':'Delhi',
        'hyderabad':'Telangana','warangal':'Telangana','mahabubnagar':'Telangana',
        'nalgonda':'Telangana','kamareddy':'Telangana','karimnagar':'Telangana',
        'shadnagar':'Telangana','khammam':'Telangana','siddipet':'Telangana',
        'chennai':'Tamil Nadu','coimbatore':'Tamil Nadu','madurai':'Tamil Nadu',
        'salem':'Tamil Nadu','trichy':'Tamil Nadu','vellore':'Tamil Nadu',
        'tirupur':'Tamil Nadu','thanjavur':'Tamil Nadu','dharmapuri':'Tamil Nadu',
        'tirunelveli':'Tamil Nadu','ramanathapuram':'Tamil Nadu','viluppuram':'Tamil Nadu',
        'kumbakonam':'Tamil Nadu','sriperumbudur':'Tamil Nadu','nagercoil':'Tamil Nadu',
        'karaikudi':'Tamil Nadu','tiruchendur':'Tamil Nadu','dharapuram':'Tamil Nadu',
        'krishnagiri':'Tamil Nadu','tada':'Tamil Nadu',
        'kolkata':'West Bengal','howrah':'West Bengal','siliguri':'West Bengal',
        'durgapur':'West Bengal','asansol':'West Bengal','midnapore':'West Bengal',
        'bagnan':'West Bengal','tarkeshwar':'West Bengal','contai':'West Bengal',
        'ahmedabad':'Gujarat','surat':'Gujarat','vadodara':'Gujarat','rajkot':'Gujarat',
        'bhavnagar':'Gujarat','jamnagar':'Gujarat','anand':'Gujarat','gandhinagar':'Gujarat',
        'kadodara':'Gujarat','bhuj':'Gujarat','morbi':'Gujarat',
        'jaipur':'Rajasthan','udaipur':'Rajasthan','jodhpur':'Rajasthan','kota':'Rajasthan',
        'bikaner':'Rajasthan','ajmer':'Rajasthan','bharatpur':'Rajasthan','sikar':'Rajasthan',
        'nadoti':'Rajasthan','ganga nagar':'Rajasthan','pali':'Rajasthan',
        'lucknow':'Uttar Pradesh','noida':'Uttar Pradesh','agra':'Uttar Pradesh',
        'kanpur':'Uttar Pradesh','varanasi':'Uttar Pradesh','ghaziabad':'Uttar Pradesh',
        'bareilly':'Uttar Pradesh','aligarh':'Uttar Pradesh','gorakhpur':'Uttar Pradesh',
        'saharanpur':'Uttar Pradesh','jhansi':'Uttar Pradesh','sultanpur':'Uttar Pradesh',
        'basti':'Uttar Pradesh','buduan':'Uttar Pradesh','deoria':'Uttar Pradesh',
        'ballia':'Uttar Pradesh','pratapgarh':'Uttar Pradesh','orai':'Uttar Pradesh',
        'banda':'Uttar Pradesh','sikandra rao':'Uttar Pradesh','malihabad':'Uttar Pradesh',
        'firozpur':'Uttar Pradesh','vrindavan':'Uttar Pradesh','najibabad':'Uttar Pradesh',
        'gurugram':'Haryana','gurgaon':'Haryana','faridabad':'Haryana','sonipat':'Haryana',
        'panipat':'Haryana','karnal':'Haryana','rohtak':'Haryana','ambala':'Haryana',
        'chandigarh':'Chandigarh',
        'patna':'Bihar','gaya':'Bihar','muzaffarpur':'Bihar','bhagalpur':'Bihar',
        'darbhanga':'Bihar','purnia':'Bihar','arrah':'Bihar','begusarai':'Bihar',
        'sasaram':'Bihar','hajipur':'Bihar','munger':'Bihar','lhakhisarai':'Bihar',
        'sitamarhi':'Bihar','gopalganj':'Bihar','samastipur':'Bihar','mohania':'Bihar',
        'motihari':'Bihar','bihar sharif':'Bihar','araria':'Bihar','khadda':'Bihar',
        'chhapra':'Bihar','aurangabad-br':'Bihar',
        'indore':'Madhya Pradesh','bhopal':'Madhya Pradesh','jabalpur':'Madhya Pradesh',
        'gwalior':'Madhya Pradesh','ujjain':'Madhya Pradesh','rewah':'Madhya Pradesh',
        'sagar':'Madhya Pradesh','ratlam':'Madhya Pradesh','balaghat':'Madhya Pradesh',
        'barwani':'Madhya Pradesh','manendragarh':'Madhya Pradesh','chhatarpur':'Madhya Pradesh',
        'raipur':'Chhattisgarh','bhilai':'Chhattisgarh','bilaspur':'Chhattisgarh',
        'dhamtari':'Chhattisgarh','durg':'Chhattisgarh','kawardha':'Chhattisgarh',
        'berla':'Chhattisgarh',
        'ranchi':'Jharkhand','jamshedpur':'Jharkhand','dhanbad':'Jharkhand',
        'giridih':'Jharkhand','deoghar':'Jharkhand','sahebganj':'Jharkhand',
        'hazaribagh':'Jharkhand','garhwa':'Jharkhand','madhupur':'Jharkhand',
        'daltonganj':'Jharkhand','pakur':'Jharkhand',
        'bhubaneswar':'Odisha','cuttack':'Odisha','sambalpur':'Odisha',
        'berhampur':'Odisha','balangir':'Odisha','baripada':'Odisha',
        'bhadrak':'Odisha','jeypore':'Odisha','behrampur':'Odisha','khurdha':'Odisha',
        'guwahati':'Assam','johrat':'Assam','silchar':'Assam','tezpur':'Assam',
        'sarthebari':'Assam','nalbari':'Assam','dibrugarh':'Assam',
        'amritsar':'Punjab','ludhiana':'Punjab','jalandhar':'Punjab','mohali':'Punjab',
        'bathinda':'Punjab','patiala':'Punjab',
        'kochi':'Kerala','ernakulam':'Kerala','thiruvananthapuram':'Kerala',
        'trivandrum':'Kerala','kozhikode':'Kerala','calicut':'Kerala',
        'thrissur':'Kerala','kannur':'Kerala','palakkad':'Kerala','kollam':'Kerala',
        'shoranur':'Kerala',
        'visakhapatnam':'Andhra Pradesh','vijayawada':'Andhra Pradesh','tirupati':'Andhra Pradesh',
        'guntur':'Andhra Pradesh','nellore':'Andhra Pradesh','kurnool':'Andhra Pradesh',
        'anantapur':'Andhra Pradesh','chittoor':'Andhra Pradesh','srikakulam':'Andhra Pradesh',
        'kakinada':'Andhra Pradesh','rajahmundry':'Andhra Pradesh','bhimavaram':'Andhra Pradesh',
        'palakollu':'Andhra Pradesh','tadepalligudem':'Andhra Pradesh','narsipatnam':'Andhra Pradesh',
        'nandigama':'Andhra Pradesh','gudivada':'Andhra Pradesh','markapur':'Andhra Pradesh',
        'samarlakota':'Andhra Pradesh','palasa':'Andhra Pradesh','narasaraopet':'Andhra Pradesh',
        'ranastalam':'Andhra Pradesh','puttur':'Andhra Pradesh','pangodu':'Andhra Pradesh',
        'dehradun':'Uttarakhand','haridwar':'Uttarakhand','rishikesh':'Uttarakhand',
        'srinagar':'Jammu & Kashmir','jammu':'Jammu & Kashmir','sopore':'Jammu & Kashmir',
        'anantnag':'Jammu & Kashmir','goa':'Goa','panaji':'Goa','margao':'Goa',
        'imphal':'Manipur','agartala':'Tripura',
        'balurghat':'West Bengal','bhadohi':'Uttar Pradesh','buduan':'Uttar Pradesh',
        'cuddapah':'Andhra Pradesh','lakhisarai':'Bihar','gaya':'Bihar',
        'dullahpur':'Uttar Pradesh','itava':'Uttar Pradesh',
        'andhra pradesh':'Andhra Pradesh','arunachal':'Arunachal Pradesh',
        'himachal':'Himachal Pradesh','uttarakhand':'Uttarakhand','meghalaya':'Meghalaya',
        'mizoram':'Mizoram','nagaland':'Nagaland','sikkim':'Sikkim',
        'tamil nadu':'Tamil Nadu','uttar pradesh':'Uttar Pradesh',
        'west bengal':'West Bengal','madhya pradesh':'Madhya Pradesh',
        'jammu & kashmir':'Jammu & Kashmir','jammu and kashmir':'Jammu & Kashmir',
        'dadra':'Dadra & Nagar Haveli','daman':'Daman & Diu',
        'puducherry':'Puducherry','pondicherry':'Puducherry',
        'andaman':'Andaman & Nicobar','lakshadweep':'Lakshadweep',
    }
    for k, v in sorted(m.items(), key=lambda x: -len(x[0])):
        if k in c:
            return v
    return city.title() or 'Unknown'

def safe_filename(name):
    return re.sub(r'[\\/:*?"<>|]+', '-', str(name or 'Report')).strip().replace(' ', '_')[:60]

# ── Report Builder ───────────────────────────────────────────────────

def _auto_col(raw, patterns):
    for p in patterns:
        for k in raw:
            if p in k:
                return str(raw.get(k, '')).strip()
    return ''

def build_report(records):
    normalized = []
    for r in records:
        raw = {k.lower().strip(): v for k, v in r.items()}
        n = {}
        id_raw = str(raw.get('order id', raw.get('reference number', raw.get('lr no', '')))).replace('`', '').strip()
        n['id'] = id_raw or _auto_col(raw, ['order', 'reference', 'lr no', 'lrn', 'tracking', 'awb', 'waybill', 'shipment'])
        client_raw = str(raw.get('client', raw.get('customer code', raw.get('customer', '')))).replace('`', '').strip()
        n['client'] = client_raw or _auto_col(raw, ['client', 'customer', 'network', 'partner', 'vendor'])
        n['boxes'] = parse_num(raw.get('no of boxes', raw.get('num pieces', _auto_col(raw, ['box', 'piece', 'item', 'quantity', 'pkg'])))) or 1
        n['origin'] = str(raw.get('origin city', raw.get('sender city', _auto_col(raw, ['origin', 'sender', 'from city', 'source'])))).strip().title()
        n['dest'] = str(raw.get('destination city', raw.get('consignee city', _auto_col(raw, ['destination', 'consignee', 'to city'])))).strip().title()
        wh = str(raw.get('client location/warehouse', raw.get('origin hub name', _auto_col(raw, ['warehouse', 'hub', 'location'])))).strip()
        n['warehouse'] = wh.title() if wh else 'Unassigned'
        n['manifest'] = parse_dt(raw.get('manifest date', raw.get('created at',
            _auto_col(raw, ['manifest', 'created', 'booking']))))
        n['pickup'] = parse_dt(raw.get('pickup date', raw.get('last pickup completed time',
            _auto_col(raw, ['pickup', 'dispatch', 'shipped']))))
        n['promise'] = parse_dt(raw.get('promise date', raw.get('expected delivery date', raw.get('expected date',
            _auto_col(raw, ['promise', 'expected', 'committed', 'scheduled'])))))
        n['delivered'] = parse_dt(raw.get('delivered date', raw.get('delivered time',
            _auto_col(raw, ['delivered', 'delivery date', 'completion']))))
        status = str(raw.get('current status', raw.get('status',
            _auto_col(raw, ['status', 'condition'])))).strip()
        n['status'] = status
        n['is_delivered'] = status.lower() in ('delivered', 'delivered to consignee') or n['delivered'] is not None
        state = str(raw.get('state', _auto_col(raw, ['state', 'province']))).strip().title()
        n['state'] = state if state else guess_state(n['dest'])
        n['amount'] = parse_num(raw.get('package amount', raw.get('declared value',
            _auto_col(raw, ['amount', 'value', 'declared', 'price']))))
        n['weight'] = parse_num(raw.get('weight', _auto_col(raw, ['weight', 'wt', 'kg'])))
        n['attempts'] = int(parse_num(raw.get('attempt count', _auto_col(raw, ['attempt', 'retry']))))
        n['last_scan'] = parse_dt(raw.get('last scan date', _auto_col(raw, ['scan', 'last update'])))
        n['remarks'] = str(raw.get('remarks', raw.get('delivery failure reason',
            _auto_col(raw, ['remarks', 'reason', 'failure', 'comment'])))).strip()
        normalized.append(n)

    shipments = len(normalized)
    delivered = sum(1 for r in normalized if r['is_delivered'])
    open_s = shipments - delivered
    late = sum(1 for r in normalized if r['is_delivered'] and r['promise'] and r['delivered'] and r['delivered'] > r['promise'])
    open_delay = sum(1 for r in normalized if not r['is_delivered'] and r['promise'] and r['promise'] < datetime.now())
    dwp = [r for r in normalized if r['is_delivered'] and r['delivered'] and r['promise']]
    on_time = sum(1 for r in dwp if r['delivered'] <= r['promise'])
    on_rate = round(on_time / len(dwp) * 100) if dwp else 0
    tat_r = [r for r in normalized if r['is_delivered'] and r['delivered'] and (r['pickup'] or r['manifest'])]
    has_del_dates = any(r['delivered'] for r in normalized)
    avg_tat = sum((r['delivered'] - (r['pickup'] or r['manifest'])).days for r in tat_r) / len(tat_r) if tat_r else 0
    del_rate = round(delivered / shipments * 100) if shipments else 0

    sc = {}
    for r in normalized:
        s = r['status'] or 'Unknown'
        sc[s] = sc.get(s, 0) + 1
    state_c = {}
    for r in normalized:
        st = r['state'] or 'Unknown'
        state_c[st] = state_c.get(st, 0) + 1
    lc = {}
    for r in normalized:
        lane = f"{r['origin']} -> {r['dest']}"
        lc[lane] = lc.get(lane, 0) + 1
    aging = {'0-2d': 0, '3-5d': 0, '6-8d': 0, '9d+': 0}
    now = datetime.now()
    for r in normalized:
        if not r['is_delivered']:
            start = r['pickup'] or r['manifest'] or r['last_scan']
            if start:
                age = max(0, (now - start).days)
                if age <= 2: aging['0-2d'] += 1
                elif age <= 5: aging['3-5d'] += 1
                elif age <= 8: aging['6-8d'] += 1
                else: aging['9d+'] += 1

    priorities = sorted(normalized, key=lambda r: (
        0 if (not r['is_delivered'] and r['promise'] and r['promise'] < now) else 1 if not r['is_delivered'] else 2,
        r['promise'] if r['promise'] else datetime.max))

    ts = sorted(sc.items(), key=lambda x: -x[1])[:6]
    tstate = sorted(state_c.items(), key=lambda x: -x[1])[:8]
    tl = sorted(lc.items(), key=lambda x: -x[1])[:10]
    top_s = tstate[0][0] if tstate else '-'
    top_l = tl[0][0] if tl else '-'
    ob = max(aging, key=aging.get)

    return {
        'client': 'Customer',
        'period': 'Selected period',
        'generated': f'Generated {datetime.now().strftime("%d %b %Y %H:%M")}',
        'shipments': shipments, 'delivered': delivered, 'open': open_s,
        'late_delivered': late, 'open_delayed': open_delay,
        'delivered_rate': del_rate, 'on_time_rate': on_rate,
        'avg_tat': avg_tat, 'avg_tat_label': f'{avg_tat:.1f}d' if (has_del_dates and tat_r) else 'N/A',
        'total_boxes': round(sum(r['boxes'] for r in normalized)),
        'total_weight': sum(r['weight'] for r in normalized),
        'total_amount': sum(r['amount'] for r in normalized),
        'attempted': sum(1 for r in normalized if r['attempts'] > 0),
        'value_label': fmt_money(sum(r['amount'] for r in normalized)),
        'weight_label': f'{sum(r["weight"] for r in normalized):.1f} kg',
        'status_entries': ts, 'state_entries': tstate, 'lane_entries': tl,
        'aging': aging, 'priority_records': priorities[:14],
        'insights': [
            f'{del_rate}% closure across {fmt_num(shipments)} shipments; {fmt_num(open_s)} remain open.',
            f'{on_rate}% on-time delivery; {fmt_num(late)} late delivered.',
            f'{fmt_num(open_delay)} open past promise - require ETA updates.',
            f'{top_s} is top destination state.', f'Top lane: {top_l}.',
        ], 'records': normalized,
    }

# ══════════════════════════════════════════════════════════════════════
#  PPT Generator
# ══════════════════════════════════════════════════════════════════════

def generate_ppt(report):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    sw, sh = Inches(13.33), Inches(7.5)
    M = Inches(0.5)

    def rect(sl, x, y, w, h, fill, line=None):
        s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = fill
        if line:
            s.line.color.rgb = line; s.line.width = Pt(0.7)
        else: s.line.fill.background()
        return s

    def txt(sl, x, y, w, h, text, size=12, bold=False, color=DARK, align=PP_ALIGN.LEFT):
        tb = sl.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = str(text)
        p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color; p.font.name = 'Aptos'; p.alignment = align
        return tb

    def bullets(sl, x, y, w, h, lines, size=10, color=DARK):
        tb = sl.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        for i, l in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(l); p.font.size = Pt(size); p.font.color.rgb = color; p.font.name = 'Aptos'; p.space_after = Pt(4)

    def kpi(sl, metrics, y):
        cw, ch, gap = Inches(2.82), Inches(0.98), Inches(0.3)
        fills = [RGBColor(0xEB,0xF5,0xFF), RGBColor(0xD1,0xFA,0xE5), RGBColor(0xFE,0xF3,0xC7), RGBColor(0xFE,0xE2,0xE2)]
        borders = [RGBColor(0x93,0xC5,0xFD), RGBColor(0x6E,0xE7,0xB7), RGBColor(0xFC,0xD3,0x4D), RGBColor(0xFC,0xA5,0xA5)]
        for i, (label, val, sub) in enumerate(metrics):
            x = M + i * (cw + gap)
            rect(sl, x, y, cw, ch, fills[i % 4], borders[i % 4])
            txt(sl, x+Inches(0.12), y+Inches(0.08), Inches(2.58), Inches(0.3), str(val), 18, True, DARK)
            txt(sl, x+Inches(0.12), y+Inches(0.46), Inches(2.58), Inches(0.17), label, 7.5, True, RGBColor(0x47,0x55,0x69))
            txt(sl, x+Inches(0.12), y+Inches(0.68), Inches(2.58), Inches(0.18), sub or '', 7, False, GRAY)

    def tbl(sl, rows, x, y, w, h, cw):
        rh = h / len(rows)
        tw = sum(cw)
        for ri, row in enumerate(rows):
            cx = x
            for ci, cv in enumerate(row):
                cw_i = int(w * cw[ci] / tw)
                is_h = ri == 0
                fill = NAVY if is_h else (WHITE if ri % 2 == 0 else RGBColor(0xF1,0xF5,0xF9))
                rect(sl, cx, y+ri*rh, cw_i, rh, fill, NAVY if is_h else RGBColor(0xDD,0xE7,0xF0))
                txt(sl, cx+Inches(0.04), y+ri*rh+Inches(0.03), cw_i-Inches(0.08), rh-Inches(0.06),
                    str(cv or ''), 7 if is_h else 6.5, is_h, WHITE if is_h else RGBColor(0x1F,0x29,0x37))
                cx += cw_i

    def chart_ph(sl, x, y, w, h, label):
        pass

    def add_chart(sl, chart_type, categories, values, x, y, w, h, title, has_legend=True):
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
        chart_data = CategoryChartData()
        chart_data.categories = [str(c) for c in categories]
        chart_data.add_series('', [int(v) for v in values])
        chart_frame = sl.shapes.add_chart(chart_type, x, y, w, h, chart_data)
        chart = chart_frame.chart
        chart.has_legend = has_legend
        if has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        chart.font.size = Pt(8)
        return chart

    def add_pie_chart(sl, cats, vals, x, y, w, h, title):
        if not cats or len(cats) < 1: return
        c = add_chart(sl, XL_CHART_TYPE.PIE, cats, vals, x, y, w, h, title, True)
        plot = c.plots[0]
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(8)
        plot.data_labels.show_category_name = True
        plot.data_labels.show_percentage = True
        plot.data_labels.show_value = False

    def add_bar_chart(sl, cats, vals, x, y, w, h, title):
        if not cats: return
        c = add_chart(sl, XL_CHART_TYPE.BAR_CLUSTERED, cats, vals, x, y, w, h, title, False)
        plot = c.plots[0]
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(8)
        plot.data_labels.show_value = True

    def add_col_chart(sl, cats, vals, x, y, w, h, title):
        if not cats: return
        c = add_chart(sl, XL_CHART_TYPE.COLUMN_CLUSTERED, cats, vals, x, y, w, h, title, False)
        plot = c.plots[0]
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(8)
        plot.data_labels.show_value = True

    def add_line_chart(sl, cats, vals, x, y, w, h, title):
        if not cats: return
        c = add_chart(sl, XL_CHART_TYPE.LINE, cats, vals, x, y, w, h, title, False)
        plot = c.plots[0]
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(8)
        plot.data_labels.show_value = True

    def header(sl, title, sub='', desc=''):
        rect(sl, 0, 0, sw, Inches(0.08), NAVY)
        txt(sl, M, Inches(0.18), Inches(7.8), Inches(0.34), title, 20, True, DARK)
        if sub: txt(sl, M, Inches(0.56), Inches(8.8), Inches(0.2), sub, 9, False, GRAY)
        if desc: txt(sl, M, Inches(0.78), Inches(8.8), Inches(0.17), desc, 8, False, LGRAY)

    def footer(sl, text, client):
        rect(sl, 0, Inches(6.92), sw, Inches(0.01), RGBColor(0xCB,0xD5,0xE1))
        txt(sl, M, Inches(6.95), Inches(7.5), Inches(0.16), text or '', 7, False, LGRAY)
        txt(sl, Inches(8), Inches(6.95), Inches(4.9), Inches(0.16), '\u00a9 2026 Sonic Business Solutions \u2022 Confidential', 7, False, LGRAY, PP_ALIGN.RIGHT)

    def callout(sl, x, y, w, h, title, items, accent=BLUE):
        rect(sl, x, y, w, h, RGBColor(0xF8,0xFA,0xFC), RGBColor(0xDD,0xE7,0xF0))
        rect(sl, x, y, Inches(0.06), h, accent, accent)
        txt(sl, x+Inches(0.2), y+Inches(0.12), w-Inches(0.36), Inches(0.2), title, 9, True, DARK)
        bullets(sl, x+Inches(0.24), y+Inches(0.4), w-Inches(0.48), h-Inches(0.5),
                [f'- {b}' for b in items[:5]], 8.2, RGBColor(0x33,0x41,0x55))

    # ── Slide 1: Cover ──
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, sw, Inches(0.55), NAVY)
    txt(sl, Inches(0.5), Inches(0.08), Inches(12.3), Inches(0.4), 'SONIC BUSINESS SOLUTIONS', 10, True, WHITE)
    rect(sl, M, Inches(1.0), Inches(12.1), Inches(0.03), BLUE)
    txt(sl, M, Inches(1.28), Inches(8.5), Inches(0.55), 'MIS Dashboard', 28, True, DARK)
    txt(sl, M+Inches(0.02), Inches(1.9), Inches(7.5), Inches(0.32), report['client'], 16, True, NAVY)
    txt(sl, M+Inches(0.02), Inches(2.28), Inches(7.5), Inches(0.25), report['period'], 11, False, GRAY)
    rect(sl, M, Inches(2.75), Inches(12.1), Inches(0.02), RGBColor(0xCB,0xD5,0xE1))
    kpi(sl, [('Total Shipments', fmt_num(report['shipments']), 'All consignments'),
             ('Delivered', f"{report['delivered_rate']}%", f"{fmt_num(report['delivered'])} closed"),
             ('On-Time Rate', f"{report['on_time_rate']}%", f"{fmt_num(report['late_delivered'])} late"),
             ('Open Delayed', fmt_num(report['open_delayed']), f"{fmt_num(report['open'])} open")], Inches(3.15))
    callout(sl, Inches(0.72), Inches(4.65), Inches(11.9), Inches(1.55), 'Executive Summary', report['insights'][:3])
    txt(sl, M, Inches(6.94), Inches(10.8), Inches(0.18), report['generated'], 7.5, False, LGRAY)
    footer(sl, 'Confidential', report['client'])

    # ── Slide 2: Exec ──
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    header(sl, 'Executive Summary', report['period'], 'Key metrics, trends, and aging')
    kpi(sl, [('Package Value', report['value_label'], 'Declared invoice value'),
             ('Total Boxes', fmt_num(report['total_boxes']), report['weight_label']),
             ('Avg TAT', report['avg_tat_label'], 'Pickup to delivery'),
             ('Attempted', fmt_num(report['attempted']), 'With attempt')], Inches(0.95))
    lx, rx, tw = M+Inches(0.05), Inches(6.75), Inches(5.75)
    callout(sl, lx, Inches(2.22), tw, Inches(2.15), 'Key Insights', report['insights'])
    sr = [['Status', 'Count', 'Share', 'Comment']]
    ts = max(1, report['shipments'])
    for s, c in report['status_entries']:
        sh = round(c / ts * 100); cm = 'Closed' if s.lower() == 'delivered' else ('Moving' if 'transit' in s.lower() else 'Monitor')
        sr.append([s, fmt_num(c), f'{sh}%', cm])
    tbl(sl, sr, rx, Inches(2.22), tw, Inches(2.15), [2.4, 1.05, 1.05, 1.25])
    # Real trend chart
    trend_data = {}
    for r in report['records']:
        d = r['manifest'] or r['pickup']
        if d:
            k = d.strftime('%d %b')
            trend_data[k] = trend_data.get(k, 0) + 1
    if trend_data:
        sd = sorted(trend_data.keys())
        add_line_chart(sl, sd[-10:], [trend_data[k] for k in sd[-10:]],
                       lx, Inches(4.7), tw, Inches(1.9), 'Daily Manifest Trend')
    else:
        add_col_chart(sl, ['No data'], [0], lx, Inches(4.7), tw, Inches(1.9), 'Daily Trend')
    ar = [['Open Aging', 'Count', 'Interpretation']]
    for lb, vl in sorted(report['aging'].items()):
        ar.append([lb, fmt_num(vl), 'Escalate' if lb == '9d+' else ('Watchlist' if lb == '6-8d' else 'Follow-up')])
    tbl(sl, ar, rx, Inches(4.82), tw, Inches(1.75), [2.2, 1.2, 2.35])
    footer(sl, report['generated'], report['client'])

    # ── Slide 3: Performance ──
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    header(sl, 'Service Performance', 'Delivery closure, promise adherence', 'Operational metrics')
    status_cats = [s for s, c in report['status_entries']]
    status_vals = [c for s, c in report['status_entries']]
    if len(status_cats) >= 2:
        add_pie_chart(sl, status_cats, status_vals, Inches(0.5), Inches(1.02),
                      Inches(5.95), Inches(4.35), 'Shipment Status Distribution')
    else:
        add_col_chart(sl, status_cats or ['No Data'], status_vals or [0],
                      Inches(0.5), Inches(1.02), Inches(5.95), Inches(4.35), 'Shipment Status')
    trend_data = {}
    for r in report['records']:
        d = r['manifest'] or r['pickup']
        if d:
            k = d.strftime('%d %b')
            trend_data[k] = trend_data.get(k, 0) + 1
    if trend_data:
        sd = sorted(trend_data.keys())
        add_line_chart(sl, sd[-10:], [trend_data[k] for k in sd[-10:]],
                       Inches(6.8), Inches(1.02), Inches(5.95), Inches(2.35), 'Daily Manifest Volume')
    else:
        add_bar_chart(sl, ['No data'], [0], Inches(6.8), Inches(1.02),
                      Inches(5.95), Inches(2.35), 'Daily Volume')
    callout(sl, Inches(6.8), Inches(3.78), Inches(5.95), Inches(1.6), 'Performance Notes', [
        f"{report['delivered_rate']}% closed delivered.", f"{report['on_time_rate']}% met promise.",
        f"{fmt_num(report['open_delayed'])} past promise."], RGBColor(0xD9,0x77,0x06))
    footer(sl, report['generated'], report['client'])

    # ── Slide 4: Network ──
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    header(sl, 'Network & Movement', 'Geographic distribution, aging, top lanes', 'States, aging, and lanes')
    state_cats = [s for s, c in report['state_entries'][:8]]
    state_vals = [c for s, c in report['state_entries'][:8]]
    if state_cats:
        add_bar_chart(sl, state_cats, state_vals, Inches(0.5), Inches(1.02),
                      Inches(5.95), Inches(3.25), 'Top Destination States')
    aging_cats = sorted(report['aging'].keys())
    aging_vals = [report['aging'][k] for k in aging_cats]
    add_col_chart(sl, aging_cats, aging_vals, Inches(6.8), Inches(1.02),
                  Inches(5.95), Inches(3.25), 'Open Shipment Aging')
    lr = [['Lane', 'Shipments', 'Delivered', 'Open', 'Rate']]
    for lane, count in report['lane_entries'][:6]:
        lr_recs = [r for r in report['records'] if f"{r['origin']} -> {r['dest']}" == lane]
        ld = sum(1 for r in lr_recs if r['is_delivered'])
        rt = round(ld / count * 100) if count else 0
        lr.append([lane[:46], fmt_num(count), fmt_num(ld), fmt_num(count-ld), f'{rt}%'])
    tbl(sl, lr, Inches(0.55), Inches(4.82), Inches(12.1), Inches(1.65), [5.3, 1.65, 1.65, 1.65, 1.85])
    footer(sl, report['generated'], report['client'])

    # ── Slide 5: Exceptions ──
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    header(sl, 'Exceptions & Priority Items', 'Delayed and open shipments', 'Actionable items')
    callout(sl, Inches(0.55), Inches(0.98), Inches(12.1), Inches(0.88), 'Recommended Actions', [
        'Prioritize open delayed — confirm ETAs immediately.',
        'Review late-delivered lanes for recurring issues.'], RGBColor(0xDC,0x26,0x26))
    er = [['Order', 'Destination', 'Status', 'Promise', 'Last Scan', 'Amount', 'Action']]
    now = datetime.now()
    for r in report['priority_records'][:8]:
        delayed = not r['is_delivered'] and r['promise'] and r['promise'] < now
        if not r['is_delivered'] or delayed:
            slbl = 'Delayed' if delayed else ('Delivered' if r['is_delivered'] else 'Open')
            act = 'Expedite / confirm ETA' if delayed else ('Review late reason' if r['is_delivered'] else 'Track')
            er.append([r['id'][:22], f"{r['dest']}, {r['state']}"[:28], slbl,
                       fmt_date(r['promise']), fmt_date(r['last_scan']), fmt_money(r['amount']), act])
    tbl(sl, er, Inches(0.55), Inches(2.18), Inches(12.1), Inches(4.2), [2.05, 2.4, 1.35, 1.35, 1.65, 1.45, 1.85])
    footer(sl, report['generated'], report['client'])

    buf = io.BytesIO(); prs.save(buf); return buf.getvalue()

# ══════════════════════════════════════════════════════════════════════
#  DOCX Generator
# ══════════════════════════════════════════════════════════════════════

def generate_docx(report):
    doc = Document()
    style = doc.styles['Normal']; style.font.name = 'Aptos'; style.font.size = DocPt(10); style.paragraph_format.space_after = DocPt(4)

    def add_heading(text, level=0, color='1E3A5F'):
        p = doc.add_heading(text, level=level)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER if level == 0 else WD_ALIGN_PARAGRAPH.LEFT
        for r in p.runs: r.font.color.rgb = DocRGB(*bytes.fromhex(color))

    add_heading('MIS DASHBOARD \u2014 SONIC BUSINESS SOLUTIONS', 0)
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run('Andheri (East), Mumbai 400069, Maharashtra, India'); r.font.size = DocPt(8); r.font.color.rgb = DocRGB(0x64,0x74,0x8B)
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(report['client']); r.bold = True; r.font.size = DocPt(18); r.font.color.rgb = DocRGB(0x25,0x63,0xEB)
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(report['period']); r.font.size = DocPt(10); r.font.color.rgb = DocRGB(0x64,0x74,0x8B)
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(report['generated']); r.font.size = DocPt(8); r.font.color.rgb = DocRGB(0x94,0xA3,0xB8)
    doc.add_paragraph().add_run().add_break()

    doc.add_heading('Performance Summary', level=1)
    kpi_data = [
        f"Total Shipments: {fmt_num(report['shipments'])}  |  Delivered: {report['delivered_rate']}%  |  On-Time: {report['on_time_rate']}%  |  Open: {fmt_num(report['open'])}",
        f"Package Value: {report['value_label']}  |  Total Boxes: {fmt_num(report['total_boxes'])}  |  Weight: {report['weight_label']}  |  Avg TAT: {report['avg_tat_label']}",
        f"Attempted: {fmt_num(report['attempted'])}  |  Late Delivered: {fmt_num(report['late_delivered'])}  |  Open Delayed: {fmt_num(report['open_delayed'])}"
    ]
    for kd in kpi_data:
        p = doc.add_paragraph(); r = p.add_run(kd); r.font.size = DocPt(10); r.font.color.rgb = DocRGB(0x1F,0x29,0x37)

    doc.add_heading('Executive Summary', level=1)
    for ins in report['insights']:
        p = doc.add_paragraph(style='List Bullet'); r = p.add_run(ins); r.font.size = DocPt(9); r.font.color.rgb = DocRGB(0x47,0x55,0x69)

    def build_table(title, headers, data_rows):
        doc.add_heading(title, level=2)
        t = doc.add_table(rows=1+len(data_rows), cols=len(headers))
        t.style = 'Table Grid'; t.alignment = WD_TABLE_ALIGNMENT.CENTER
        for i, h in enumerate(headers):
            cell = t.rows[0].cells[i]; cell.text = h
            for p in cell.paragraphs:
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for r in p.runs: r.bold = True; r.font.size = DocPt(9); r.font.color.rgb = DocRGB(0x1E,0x3A,0x5F)
            from docx.oxml.ns import qn
            sh = cell._element.get_or_add_tcPr()
            sh.append(sh.makeelement(qn('w:shd'), {qn('w:fill'): 'EBF5FF', qn('w:val'): 'clear'}))
        for ri, row in enumerate(data_rows):
            for ci, cv in enumerate(row):
                t.rows[ri+1].cells[ci].text = str(cv)
        doc.add_paragraph()

    build_table('Status Breakdown', ['Status', 'Count', 'Share', 'Comment'],
                [[s, fmt_num(c), f"{round(c/max(1,report['shipments'])*100)}%",
                  'Closed' if s.lower()=='delivered' else ('Moving' if 'transit' in s.lower() else 'Monitor')]
                 for s, c in report['status_entries']])
    build_table('Open Aging', ['Aging', 'Count', 'Interpretation'],
                [[lb, fmt_num(vl), 'Escalate' if lb=='9d+' else ('Watchlist' if lb=='6-8d' else 'Follow-up')]
                 for lb, vl in sorted(report['aging'].items())])
    lane_data = []
    for la, co in report['lane_entries'][:6]:
        ld = sum(1 for r in report['records']
                 if f"{r['origin']} -> {r['dest']}" == la and r['is_delivered'])
        lo = co - ld
        lr = round(ld / max(1, co) * 100)
        lane_data.append([la[:50], fmt_num(co), fmt_num(ld), fmt_num(lo), f'{lr}%'])
    build_table('Top Lanes', ['Lane', 'Shipments', 'Delivered', 'Open', 'Rate'], lane_data)

    now = datetime.now()
    exc = []
    for r in report['priority_records'][:8]:
        d = not r['is_delivered'] and r['promise'] and r['promise'] < now
        if not r['is_delivered'] or d:
            exc.append([r['id'][:22], f"{r['dest']}, {r['state']}"[:28], 'Delayed' if d else ('Delivered' if r['is_delivered'] else 'Open'),
                        fmt_date(r['promise']), fmt_date(r['last_scan']), fmt_money(r['amount']),
                        'Expedite' if d else ('Review' if r['is_delivered'] else 'Track')])
    if exc:
        build_table('Priority Exceptions', ['Order', 'Destination', 'Status', 'Promise', 'Last Scan', 'Amount', 'Action'], exc)

    buf = io.BytesIO(); doc.save(buf); return buf.getvalue()

# ══════════════════════════════════════════════════════════════════════
#  PDF Generator (fpdf2)
# ══════════════════════════════════════════════════════════════════════

class LogisticsPDF(FPDF):
    def header(self):
        self.set_font('Helvetica', 'B', 8)
        self.set_text_color(0x1E, 0x3A, 0x5F)
        self.cell(0, 6, 'SONIC BUSINESS SOLUTIONS - MIS DASHBOARD', align='C'); self.ln(4)
        self.set_font('Helvetica', '', 6)
        self.set_text_color(0x94, 0xA3, 0xB8)
        self.cell(0, 4, 'Andheri (East), Mumbai 400069, Maharashtra, India', align='C'); self.ln(8)
        self.set_draw_color(0x25, 0x63, 0xEB); self.set_line_width(0.5)
        self.line(10, self.get_y(), 200, self.get_y()); self.ln(4)
    def footer(self):
        self.set_y(-15); self.set_font('Helvetica', '', 7)
        self.set_text_color(0x94, 0xA3, 0xB8)
        self.cell(0, 10, f'Page {self.page_no()}/{{nb}}', align='C')
    def section(self, title):
        self.set_font('Helvetica', 'B', 13); self.set_text_color(0x1E,0x3A,0x5F)
        self.cell(0, 8, title); self.ln(6)
        self.set_draw_color(0x25,0x63,0xEB); self.set_line_width(0.3)
        self.line(10, self.get_y(), 200, self.get_y()); self.ln(4)
    def sub(self, title):
        self.set_font('Helvetica', 'B', 10); self.set_text_color(0x1E,0x3A,0x5F)
        self.cell(0, 7, title); self.ln(5)
    def body(self, text, size=9, color=(0x33,0x33,0x33)):
        self.set_font('Helvetica', '', size); self.set_text_color(*color)
        self.multi_cell(0, 5, safe_text(text)); self.ln(1)
    def kpi(self, x, y, w, h, val, label, sub, bg, border):
        self.set_fill_color(*bg); self.set_draw_color(*border); self.set_line_width(0.3)
        self.rect(x, y, w, h, 'DF')
        self.set_xy(x+2, y+1.5); self.set_font('Helvetica','B',14); self.set_text_color(0x0F,0x17,0x2A)
        self.cell(w-4, 6, safe_text(str(val)))
        self.set_xy(x+2, y+8); self.set_font('Helvetica','B',6.5); self.set_text_color(0x47,0x55,0x69)
        self.cell(w-4, 4, safe_text(label))
        self.set_xy(x+2, y+12.5); self.set_font('Helvetica','',6); self.set_text_color(0x64,0x74,0x8B)
        self.cell(w-4, 4, safe_text(sub))
    def table(self, headers, rows, cw=None, hbg=(0x1E,0x3A,0x5F)):
        if not rows: return
        if not cw: cw = [190//len(headers)]*len(headers)
        tw = sum(cw); cw = [w*190/tw for w in cw]
        self.set_fill_color(*hbg); self.set_text_color(255,255,255); self.set_font('Helvetica','B',6.5)
        for i, h in enumerate(headers):
            self.cell(cw[i], 6, h, border=1, fill=True, align='C')
        self.ln()
        for ri, row in enumerate(rows):
            self.set_fill_color(255,255,255 if ri%2==0 else 0xF8) if ri%2==0 else self.set_fill_color(0xF8,0xFA,0xFC)
            self.set_text_color(0x33,0x33,0x33); self.set_font('Helvetica','',6.5)
            for ci, cv in enumerate(row):
                self.cell(cw[ci], 5, safe_text(str(cv or ''))[:60], border=1, fill=True)
            self.ln()

def generate_pdf(report):
    pdf = LogisticsPDF(); pdf.alias_nb_pages(); pdf.set_auto_page_break(auto=True, margin=20); pdf.add_page()
    pdf.set_font('Helvetica', 'B', 20); pdf.set_text_color(0x1E,0x3A,0x5F)
    pdf.cell(0, 10, 'MIS Dashboard', align='C'); pdf.ln(12)
    pdf.set_font('Helvetica', 'B', 12); pdf.set_text_color(0x25,0x63,0xEB)
    pdf.cell(0, 7, report['client'], align='C'); pdf.ln(7)
    pdf.set_font('Helvetica', '', 8); pdf.set_text_color(0x64,0x74,0x8B)
    pdf.cell(0, 5, report['period'], align='C'); pdf.ln(5)
    pdf.set_font('Helvetica', '', 7); pdf.set_text_color(0x94,0xA3,0xB8)
    pdf.cell(0, 5, report['generated'], align='C'); pdf.ln(10)
    cards = [
        (fmt_num(report['shipments']), 'Total Shipments', 'All consignments', (0xEB,0xF5,0xFF), (0x93,0xC5,0xFD)),
        (f"{report['delivered_rate']}%", 'Delivered', f"{fmt_num(report['delivered'])} closed", (0xD1,0xFA,0xE5), (0x6E,0xE7,0xB7)),
        (f"{report['on_time_rate']}%", 'On-Time', f"{fmt_num(report['late_delivered'])} late", (0xFE,0xF3,0xC7), (0xFC,0xD3,0x4D)),
        (fmt_num(report['open_delayed']), 'Open Delayed', f"{fmt_num(report['open'])} open", (0xFE,0xE2,0xE2), (0xFC,0xA5,0xA5))]
    for i, (v, l, s, bg, brd) in enumerate(cards):
        pdf.kpi(10+i*48, pdf.get_y(), 46, 16, v, l, s, bg, brd)
    pdf.ln(20)
    pdf.section('Executive Summary')
    for ins in report['insights']: pdf.body(f'- {ins}', 8, (0x47,0x55,0x69))
    pdf.sub('Status Breakdown')
    pdf.table(['Status', 'Count', 'Share', 'Comment'],
              [[s, fmt_num(c), f"{round(c/max(1,report['shipments'])*100)}%",
                'Closed' if s.lower()=='delivered' else ('Moving' if 'transit' in s.lower() else 'Monitor')]
               for s, c in report['status_entries']]); pdf.ln(4)
    pdf.sub('Open Aging')
    pdf.table(['Aging', 'Count', 'Interpretation'],
              [[lb, fmt_num(vl), 'Escalate' if lb=='9d+' else ('Watchlist' if lb=='6-8d' else 'Follow-up')]
               for lb, vl in sorted(report['aging'].items())]); pdf.ln(4)
    pdf.sub('Top Lanes')
    pdf_lanes = []
    for la, co in report['lane_entries'][:6]:
        ld = sum(1 for r in report['records']
                 if f"{r['origin']} -> {r['dest']}" == la and r['is_delivered'])
        pdf_lanes.append([la[:50], fmt_num(co), fmt_num(ld), fmt_num(co-ld), f'{round(ld/max(1,co)*100)}%'])
    pdf.table(['Lane', 'Shipments', 'Delivered', 'Open', 'Rate'], pdf_lanes); pdf.ln(4)
    now = datetime.now(); exc = []
    for r in report['priority_records'][:8]:
        d = not r['is_delivered'] and r['promise'] and r['promise'] < now
        if not r['is_delivered'] or d:
            exc.append([r['id'][:22], f"{r['dest']}, {r['state']}"[:28], 'Delayed' if d else ('Delivered' if r['is_delivered'] else 'Open'),
                        fmt_date(r['promise']), fmt_date(r['last_scan']), fmt_money(r['amount']),
                        'Expedite' if d else ('Review' if r['is_delivered'] else 'Track')])
    if exc:
        pdf.sub('Priority Exceptions')
        pdf.table(['Order', 'Destination', 'Status', 'Promise', 'Last Scan', 'Amount', 'Action'], exc,
                  cw=[28,28,18,24,24,22,22], hbg=(0x8B,0x00,0x00))
    buf = io.BytesIO(); pdf.output(buf); return buf.getvalue()

# ── App start ────────────────────────────────────────────────────────

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=os.environ.get('FLASK_DEBUG', '0') == '1')
